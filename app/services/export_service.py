from __future__ import annotations

import json
import textwrap
from io import BytesIO
from pathlib import Path

from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import inch, cm
from reportlab.lib.colors import HexColor, white, black
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT

from app.core.config import settings
from app.utils.helpers import ensure_dir, slugify


# ═══════════════════════════════════════════════════════════════
#  COLOR PALETTE for presentations
# ═══════════════════════════════════════════════════════════════

THEME = {
    "primary": RGBColor(0x6D, 0x28, 0xD9),       # Purple
    "primary_light": RGBColor(0x8B, 0x5C, 0xF6),  # Light purple
    "accent": RGBColor(0xA7, 0x8B, 0xFA),          # Lavender
    "dark_bg": RGBColor(0x0F, 0x0F, 0x1A),         # Dark background
    "card_bg": RGBColor(0x1A, 0x1A, 0x35),          # Card background
    "text_light": RGBColor(0xF1, 0xF5, 0xF9),      # Light text
    "text_muted": RGBColor(0x94, 0xA3, 0xB8),      # Muted text
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "success": RGBColor(0x22, 0xC5, 0x5E),
    "warning": RGBColor(0xF5, 0x9E, 0x0B),
}

POSTER_COLORS = {
    "bg_dark": "#0f0f1a",
    "bg_card": "#1a1a35",
    "primary": "#7c3aed",
    "primary_light": "#8b5cf6",
    "accent": "#a78bfa",
    "text_light": "#f1f5f9",
    "text_muted": "#94a3b8",
    "success": "#22c55e",
    "warning": "#f59e0b",
    "danger": "#ef4444",
    "info": "#3b82f6",
    "section_colors": [
        "#7c3aed", "#3b82f6", "#22c55e", "#f59e0b",
        "#ef4444", "#ec4899", "#06b6d4", "#8b5cf6",
    ],
}


class ExportService:

    # ═══════════════════════════════════════════════════════════
    #  SIMPLE TEXT-TO-PDF (for reports, revision, etc.)
    # ═══════════════════════════════════════════════════════════

    def export_text_to_pdf(self, title: str, content: str) -> bytes:
        buffer = BytesIO()
        pdf = canvas.Canvas(buffer, pagesize=A4)
        width, height = A4
        y = height - 60
        pdf.setFont("Helvetica-Bold", 16)
        pdf.drawString(50, y, title[:80])
        y -= 30
        pdf.setFont("Helvetica", 10)

        for raw_line in content.splitlines():
            line = raw_line.strip() or " "
            if y < 50:
                pdf.showPage()
                pdf.setFont("Helvetica", 10)
                y = height - 50
            pdf.drawString(50, y, line[:110])
            y -= 14

        pdf.save()
        buffer.seek(0)
        return buffer.read()

    # ═══════════════════════════════════════════════════════════
    #  PROFESSIONAL PPT GENERATION
    # ═══════════════════════════════════════════════════════════

    def _set_slide_bg(self, slide, r: int, g: int, b: int):
        """Set a solid background color for a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    def _add_text_box(self, slide, left, top, width, height, text,
                      font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT,
                      font_name="Calibri"):
        """Add a text box to a slide with custom formatting."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.name = font_name
        p.font.color.rgb = color or THEME["text_light"]
        p.alignment = alignment
        return tf

    def _add_bullet_frame(self, slide, left, top, width, height, bullets,
                          font_size=16, color=None, font_name="Calibri"):
        """Add a text frame with bullet points."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Handle sub-bullets (lines starting with -)
            text = bullet.strip()
            is_sub = text.startswith("-") or text.startswith("•")
            if is_sub:
                text = text.lstrip("-•").strip()
                p.text = f"  ▸ {text}"
                p.font.size = Pt(font_size - 2)
                p.font.color.rgb = color or THEME["text_muted"]
            else:
                p.text = f"● {text}"
                p.font.size = Pt(font_size)
                p.font.color.rgb = color or THEME["text_light"]

            p.font.name = font_name
            p.space_after = Pt(6)

        return tf

    def _add_accent_bar(self, slide, left, top, width, height):
        """Add a colored accent bar/rectangle to a slide."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME["primary"]
        shape.line.fill.background()
        return shape

    def _add_slide_number(self, slide, num, total, slide_width, slide_height):
        """Add slide number to bottom right."""
        self._add_text_box(
            slide,
            slide_width - Inches(1.5),
            slide_height - Inches(0.5),
            Inches(1.2),
            Inches(0.3),
            f"{num} / {total}",
            font_size=10,
            color=THEME["text_muted"],
            alignment=PP_ALIGN.RIGHT,
        )

    def export_slide_text_to_ppt(self, title: str, content: str) -> bytes:
        """Generate a professionally styled PowerPoint presentation."""
        prs = Presentation()

        # Set slide size to widescreen 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        sw = prs.slide_width
        sh = prs.slide_height

        # Parse slides from content — prefer ---SLIDE--- delimiter, fallback to double-blank-lines
        if "---SLIDE---" in content:
            raw_blocks = [block.strip() for block in content.split("---SLIDE---") if block.strip()]
        else:
            raw_blocks = [block.strip() for block in content.split("\n\n") if block.strip()]

        # Build structured slides
        slides_data = []
        for block in raw_blocks:
            lines = block.splitlines()
            if not lines:
                continue
            slide_title = lines[0].strip().lstrip("#").strip()
            # Remove markdown formatting artifacts
            slide_title = slide_title.replace("**", "").replace("*", "").strip()
            if not slide_title:
                continue
            bullets = []
            for line in lines[1:]:
                cleaned = line.strip().lstrip("-•*").strip()
                if cleaned:
                    # Remove markdown bold/italic
                    cleaned = cleaned.replace("**", "").replace("*", "")
                    bullets.append(cleaned)
            # Ensure every slide has at least one bullet point
            if not bullets:
                bullets = [f"Key concepts related to {slide_title} — see detailed notes for more information."]
            slides_data.append({"title": slide_title, "bullets": bullets})

        if not slides_data:
            slides_data = [{"title": title, "bullets": ["Content generated by StudyBuddy AI"]}]

        total_slides = len(slides_data)

        for idx, sd in enumerate(slides_data):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Dark background
            self._set_slide_bg(slide, 0x0F, 0x0F, 0x1A)

            if idx == 0:
                # ── TITLE SLIDE ──
                # Large accent bar at top
                self._add_accent_bar(slide, Inches(0), Inches(0), sw, Inches(0.08))

                # Side accent stripe
                self._add_accent_bar(slide, Inches(0.6), Inches(1.5), Inches(0.06), Inches(2.5))

                # Title
                self._add_text_box(
                    slide, Inches(1.0), Inches(1.8), Inches(10), Inches(1.5),
                    sd["title"],
                    font_size=40, bold=True, color=THEME["white"],
                    font_name="Calibri",
                )

                # Subtitle / first bullet as subtitle
                subtitle_text = "Generated by StudyBuddy AI"
                if sd["bullets"]:
                    subtitle_text = sd["bullets"][0]
                self._add_text_box(
                    slide, Inches(1.0), Inches(3.5), Inches(10), Inches(0.8),
                    subtitle_text,
                    font_size=20, color=THEME["accent"],
                    font_name="Calibri",
                )

                # Bottom accent bar
                self._add_accent_bar(slide, Inches(0), sh - Inches(0.08), sw, Inches(0.08))

                # Branding
                self._add_text_box(
                    slide, Inches(0.8), sh - Inches(0.7), Inches(4), Inches(0.4),
                    "🎓 StudyBuddy AI",
                    font_size=12, color=THEME["text_muted"],
                )

            elif idx == total_slides - 1:
                # ── CONCLUSION SLIDE ──
                self._add_accent_bar(slide, Inches(0), Inches(0), sw, Inches(0.08))

                self._add_text_box(
                    slide, Inches(1.0), Inches(0.8), Inches(10), Inches(0.8),
                    sd["title"],
                    font_size=36, bold=True, color=THEME["white"],
                )

                # Accent line under title
                self._add_accent_bar(slide, Inches(1.0), Inches(1.7), Inches(3), Inches(0.04))

                if sd["bullets"]:
                    self._add_bullet_frame(
                        slide, Inches(1.0), Inches(2.0), Inches(11), Inches(4.5),
                        sd["bullets"], font_size=18,
                    )

                self._add_accent_bar(slide, Inches(0), sh - Inches(0.08), sw, Inches(0.08))
                self._add_slide_number(slide, idx + 1, total_slides, sw, sh)

            else:
                # ── CONTENT SLIDES ──
                # Top accent bar
                self._add_accent_bar(slide, Inches(0), Inches(0), sw, Inches(0.06))

                # Side accent stripe
                self._add_accent_bar(slide, Inches(0.5), Inches(0.6), Inches(0.05), Inches(0.8))

                # Title
                self._add_text_box(
                    slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                    sd["title"],
                    font_size=30, bold=True, color=THEME["white"],
                )

                # Divider line
                self._add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(4), Inches(0.03))

                # Bullet points
                if sd["bullets"]:
                    self._add_bullet_frame(
                        slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0),
                        sd["bullets"], font_size=18,
                    )

                # Bottom bar
                self._add_accent_bar(slide, Inches(0), sh - Inches(0.06), sw, Inches(0.06))

                # Slide number
                self._add_slide_number(slide, idx + 1, total_slides, sw, sh)

        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    # ═══════════════════════════════════════════════════════════
    #  COLORFUL POSTER PDF GENERATION
    # ═══════════════════════════════════════════════════════════

    def _draw_rounded_rect(self, c, x, y, w, h, radius, fill_color):
        """Draw a rounded rectangle on the canvas."""
        c.setFillColor(HexColor(fill_color))
        c.roundRect(x, y, w, h, radius, fill=1, stroke=0)

    def export_poster_to_pdf(self, poster_data: dict) -> bytes:
        """Generate a colorful, visually engaging poster PDF.

        poster_data should contain:
        - title: str
        - tagline: str
        - sections: list of {heading, points: list[str]}
        - conclusion: str
        """
        buffer = BytesIO()
        page_w, page_h = landscape(A4)
        c = canvas.Canvas(buffer, pagesize=landscape(A4))

        title = poster_data.get("title", "StudyBuddy Poster")
        tagline = poster_data.get("tagline", "AI-Generated Learning Poster")
        sections = poster_data.get("sections", [])
        conclusion = poster_data.get("conclusion", "")

        # ── Background ──
        c.setFillColor(HexColor("#0f0f1a"))
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)

        # ── Top Header Bar ──
        c.setFillColor(HexColor("#7c3aed"))
        c.rect(0, page_h - 120, page_w, 120, fill=1, stroke=0)

        # ── Accent stripe under header ──
        c.setFillColor(HexColor("#a78bfa"))
        c.rect(0, page_h - 125, page_w, 5, fill=1, stroke=0)

        # ── Title ──
        c.setFillColor(HexColor("#ffffff"))
        c.setFont("Helvetica-Bold", 28)
        # Center the title
        title_display = title[:60]
        title_w = c.stringWidth(title_display, "Helvetica-Bold", 28)
        c.drawString((page_w - title_w) / 2, page_h - 55, title_display)

        # ── Tagline ──
        c.setFillColor(HexColor("#e2e8f0"))
        c.setFont("Helvetica-Oblique", 14)
        tagline_display = tagline[:100]
        tagline_w = c.stringWidth(tagline_display, "Helvetica-Oblique", 14)
        c.drawString((page_w - tagline_w) / 2, page_h - 85, tagline_display)

        # ── Branding ──
        c.setFillColor(HexColor("#e2e8f0"))
        c.setFont("Helvetica", 10)
        c.drawString(30, page_h - 110, "🎓 StudyBuddy AI")

        # ── Content Sections ──
        section_colors = POSTER_COLORS["section_colors"]
        margin = 30
        y_cursor = page_h - 155

        if sections:
            # Try to lay out sections in 2 columns
            col_width = (page_w - margin * 3) / 2
            col_x = [margin, margin * 2 + col_width]
            col_y = [y_cursor, y_cursor]

            for i, section in enumerate(sections):
                col_idx = 0 if col_y[0] >= col_y[1] else 1
                x = col_x[col_idx]
                y = col_y[col_idx]

                heading = section.get("heading", f"Section {i + 1}")
                points = section.get("points", [])

                # Calculate height needed
                section_height = 35 + len(points) * 18 + 15
                if section_height > 200:
                    section_height = 200

                # Check if we need a new page
                if y - section_height < 60:
                    # Skip this section if no space
                    continue

                # Section card background
                self._draw_rounded_rect(c, x, y - section_height, col_width, section_height, 8, "#1a1a35")

                # Section color accent bar (left side)
                color = section_colors[i % len(section_colors)]
                c.setFillColor(HexColor(color))
                c.roundRect(x, y - section_height, 5, section_height, 2, fill=1, stroke=0)

                # Section heading
                c.setFillColor(HexColor(color))
                c.setFont("Helvetica-Bold", 13)
                c.drawString(x + 15, y - 22, heading[:50])

                # Section points
                c.setFillColor(HexColor("#e2e8f0"))
                c.setFont("Helvetica", 10)
                point_y = y - 42
                for point in points[:8]:  # max 8 points per section
                    if point_y < y - section_height + 10:
                        break
                    point_text = point[:80]
                    c.drawString(x + 20, point_y, f"▸ {point_text}")
                    point_y -= 16

                col_y[col_idx] = y - section_height - 12

        # ── Conclusion bar at bottom ──
        if conclusion:
            c.setFillColor(HexColor("#1a1a35"))
            c.roundRect(margin, 20, page_w - margin * 2, 40, 8, fill=1, stroke=0)

            # Accent
            c.setFillColor(HexColor("#22c55e"))
            c.roundRect(margin, 20, 5, 40, 2, fill=1, stroke=0)

            c.setFillColor(HexColor("#e2e8f0"))
            c.setFont("Helvetica-Bold", 11)
            c.drawString(margin + 15, 43, f"💡 {conclusion[:120]}")

        c.save()
        buffer.seek(0)
        return buffer.read()

    # ═══════════════════════════════════════════════════════════
    #  MIND MAP HTML (Mermaid diagram)
    # ═══════════════════════════════════════════════════════════

    def export_mindmap_to_html(self, mermaid_code: str) -> str:
        """Generate a self-contained HTML page with a Mermaid mind map diagram and download button."""
        return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>StudyBuddy Mind Map</title>
<script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
<style>
  body {{
    background: #0f0f1a;
    color: #e2e8f0;
    font-family: 'Inter', 'Segoe UI', sans-serif;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    min-height: 100vh;
    margin: 0;
    padding: 20px;
  }}
  h1 {{
    color: #a78bfa;
    font-size: 1.5rem;
    margin-bottom: 1rem;
  }}
  .mermaid {{
    background: #1a1a35;
    border-radius: 16px;
    padding: 2rem;
    border: 1px solid rgba(124,58,237,0.2);
    max-width: 95vw;
    overflow: auto;
  }}
  .mermaid svg {{
    max-width: 100%;
  }}
  .download-btn {{
    margin-top: 1.2rem;
    padding: 0.7rem 1.8rem;
    background: linear-gradient(135deg, #6d28d9, #7c3aed, #8b5cf6);
    color: #ffffff;
    border: none;
    border-radius: 12px;
    font-size: 1rem;
    font-weight: 600;
    cursor: pointer;
    display: inline-flex;
    align-items: center;
    gap: 0.5rem;
    box-shadow: 0 4px 14px rgba(124,58,237,0.3);
    transition: all 0.25s ease;
  }}
  .download-btn:hover {{
    transform: translateY(-2px);
    box-shadow: 0 6px 20px rgba(124,58,237,0.45);
  }}
  .footer {{
    margin-top: 1rem;
    font-size: 0.8rem;
    color: #64748b;
  }}
</style>
</head>
<body>
<h1>🧠 StudyBuddy Mind Map</h1>
<div class="mermaid" id="mindmap-container">
{mermaid_code}
</div>
<button class="download-btn" onclick="downloadPNG()" id="downloadBtn">⬇️ Download as PNG</button>
<div class="footer">Generated by StudyBuddy AI</div>
<script>
  mermaid.initialize({{
    startOnLoad: true,
    theme: 'dark',
    themeVariables: {{
      primaryColor: '#7c3aed',
      primaryTextColor: '#f1f5f9',
      primaryBorderColor: '#8b5cf6',
      lineColor: '#a78bfa',
      secondaryColor: '#1a1a35',
      tertiaryColor: '#141428',
      fontSize: '14px'
    }}
  }});

  function downloadPNG() {{
    const svgEl = document.querySelector('#mindmap-container svg');
    if (!svgEl) {{ alert('Mind map not rendered yet. Please wait.'); return; }}
    const svgData = new XMLSerializer().serializeToString(svgEl);
    const canvas = document.createElement('canvas');
    const ctx = canvas.getContext('2d');
    const img = new Image();
    const svgBlob = new Blob([svgData], {{type: 'image/svg+xml;charset=utf-8'}});
    const url = URL.createObjectURL(svgBlob);
    img.onload = function() {{
      const scale = 2;
      canvas.width = img.width * scale;
      canvas.height = img.height * scale;
      ctx.fillStyle = '#0f0f1a';
      ctx.fillRect(0, 0, canvas.width, canvas.height);
      ctx.drawImage(img, 0, 0, canvas.width, canvas.height);
      URL.revokeObjectURL(url);
      canvas.toBlob(function(blob) {{
        const a = document.createElement('a');
        a.href = URL.createObjectURL(blob);
        a.download = 'studybuddy_mindmap.png';
        a.click();
        URL.revokeObjectURL(a.href);
      }}, 'image/png');
    }};
    img.src = url;
  }}
</script>
</body>
</html>"""

    # ═══════════════════════════════════════════════════════════
    #  TEXT-TO-SPEECH
    # ═══════════════════════════════════════════════════════════

    def text_to_speech_file(self, title: str, content: str) -> str:
        export_dir = ensure_dir(settings.export_dir)
        output_path = str(Path(export_dir) / f"{slugify(title)}.mp3")
        gTTS(text=content[:4000], lang="en").save(output_path)
        return output_path

    # ═══════════════════════════════════════════════════════════
    #  EXAM PREDICTIONS PDF (using fpdf2)
    # ═══════════════════════════════════════════════════════════

    def export_predictions_to_pdf(self, predictions: list[dict], subject: str = "Exam") -> bytes:
        """Generate a styled PDF document of predicted exam questions.

        Uses fpdf2 to create a professional-looking PDF with question numbers,
        confidence levels (color-coded), topics, and reasoning.
        """
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # ── Title ──
        pdf.set_fill_color(109, 40, 217)  # Purple header
        pdf.rect(0, 0, 210, 35, "F")
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Helvetica", "B", 20)
        pdf.set_y(8)
        pdf.cell(0, 12, f"Predicted Exam Questions", align="C", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 11)
        pdf.cell(0, 8, f"Subject: {subject} | Generated by StudyBuddy AI", align="C", new_x="LMARGIN", new_y="NEXT")

        pdf.set_y(42)
        pdf.set_text_color(0, 0, 0)

        # ── Questions ──
        for i, pred in enumerate(predictions):
            question = pred.get("predicted_question", pred.get("question", "N/A"))
            confidence = pred.get("confidence", "Medium")
            topic = pred.get("topic", "")
            reasoning = pred.get("reasoning", "")

            # Confidence color
            if confidence.lower() == "high":
                r, g, b = 34, 197, 94  # Green
            elif confidence.lower() == "medium":
                r, g, b = 245, 158, 11  # Amber
            else:
                r, g, b = 239, 68, 68  # Red

            # Question number + confidence badge
            pdf.set_font("Helvetica", "B", 12)
            pdf.set_text_color(50, 50, 50)
            pdf.cell(0, 8, f"Q{i + 1}. {question[:90]}", new_x="LMARGIN", new_y="NEXT")

            # Handle long questions that overflow
            if len(question) > 90:
                pdf.set_font("Helvetica", "", 10)
                pdf.cell(0, 6, f"      {question[90:180]}", new_x="LMARGIN", new_y="NEXT")

            # Confidence + Topic line
            pdf.set_font("Helvetica", "", 9)
            pdf.set_text_color(r, g, b)
            conf_line = f"   Confidence: {confidence}"
            if topic:
                conf_line += f"  |  Topic: {topic}"
            pdf.cell(0, 6, conf_line, new_x="LMARGIN", new_y="NEXT")

            # Reasoning
            if reasoning:
                pdf.set_text_color(100, 100, 100)
                pdf.set_font("Helvetica", "I", 9)
                pdf.multi_cell(0, 5, f"   Reasoning: {reasoning[:200]}")

            pdf.ln(4)

        # ── Footer ──
        pdf.set_y(-25)
        pdf.set_font("Helvetica", "I", 8)
        pdf.set_text_color(150, 150, 150)
        pdf.cell(0, 10, "Generated by StudyBuddy AI - Smart Learning Assistant", align="C")

        # Return PDF bytes
        return bytes(pdf.output())

    # ═══════════════════════════════════════════════════════════
    #  EMAIL WITH ATTACHMENT (using yagmail)
    # ═══════════════════════════════════════════════════════════

    def send_email_with_attachment(
        self, to_email: str, subject: str, body: str,
        attachment_bytes: bytes, attachment_filename: str
    ) -> bool:
        """Send an email with a PDF attachment using yagmail.

        Uses SMTP credentials from app settings (.env file).
        Returns True if email was sent successfully, False otherwise.
        """
        import yagmail
        import tempfile
        import os

        # Save attachment to temp file (yagmail needs file path)
        temp_dir = ensure_dir(settings.export_dir)
        temp_path = os.path.join(temp_dir, attachment_filename)
        try:
            with open(temp_path, "wb") as f:
                f.write(attachment_bytes)

            # Initialize yagmail with SMTP credentials from .env
            yag = yagmail.SMTP(
                user=settings.smtp_email,
                password=settings.smtp_password,
                host=settings.smtp_server,
                port=settings.smtp_port,
            )

            # Send the email
            yag.send(
                to=to_email,
                subject=subject,
                contents=body,
                attachments=temp_path,
            )
            return True

        except Exception as exc:
            raise ValueError(f"Failed to send email: {exc}")
        finally:
            # Clean up temp file
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except OSError:
                    pass

