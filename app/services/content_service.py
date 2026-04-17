from __future__ import annotations

import io

import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi

from app.utils.helpers import extract_youtube_video_id, file_extension


class ContentService:
    def extract_text_from_pdf(self, file_bytes: bytes) -> str:
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join((page.extract_text() or "") for page in reader.pages).strip()

    def extract_text_from_docx(self, file_bytes: bytes) -> str:
        document = Document(io.BytesIO(file_bytes))
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()

    def extract_text_from_txt(self, file_bytes: bytes) -> str:
        return file_bytes.decode("utf-8", errors="ignore").strip()

    def extract_text_from_ppt(self, file_bytes: bytes) -> str:
        presentation = Presentation(io.BytesIO(file_bytes))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines).strip()

    def extract_text_from_upload(self, filename: str, file_bytes: bytes) -> str:
        extension = file_extension(filename)
        if extension == ".pdf":
            return self.extract_text_from_pdf(file_bytes)
        if extension == ".docx":
            return self.extract_text_from_docx(file_bytes)
        if extension == ".txt":
            return self.extract_text_from_txt(file_bytes)
        if extension in {".ppt", ".pptx"}:
            return self.extract_text_from_ppt(file_bytes)
        raise ValueError("Unsupported file type. Please upload PDF, DOCX, TXT, or PPT/PPTX.")

    def extract_youtube_transcript(self, youtube_url: str) -> str:
        video_id = extract_youtube_video_id(youtube_url)
        if not video_id:
            raise ValueError("Invalid YouTube URL.")
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        return " ".join(item["text"] for item in transcript).strip()

    def extract_website_content(self, url: str) -> str:
        response = requests.get(url, timeout=20, headers={"User-Agent": "StudyBuddy/1.0"})
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return " ".join(soup.stripped_strings).strip()

    def extract_pdf_title(self, file_bytes: bytes) -> str:
        """Extract the heading/title from the first page of a PDF.

        Reads the first page text and returns the first non-empty line
        as the document title. Falls back to 'Untitled Document' if
        no text is found.
        """
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            if reader.pages:
                first_page_text = reader.pages[0].extract_text() or ""
                for line in first_page_text.splitlines():
                    cleaned = line.strip()
                    if cleaned and len(cleaned) > 2:
                        return cleaned[:120]  # First meaningful line = title
            return "Untitled Document"
        except Exception:
            return "Untitled Document"

